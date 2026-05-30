"""
Testes para core/pptx_converter.py.
Fixtures criadas programaticamente via python-pptx (sem arquivos binários no repo).
"""
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from core.pptx_converter import _tabela_para_md, pptx_to_md

# ── Fixtures ────────────────────────────────────────────────────────────────

@pytest.fixture
def pptx_basico(tmp_path) -> Path:
    """PPTX com 2 slides: texto simples."""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # título + conteúdo

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Slide Um"
    slide1.placeholders[1].text = "Conteúdo do slide um."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Slide Dois"
    slide2.placeholders[1].text = "Conteúdo do slide dois."

    path = tmp_path / "basico.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def pptx_com_tabela(tmp_path) -> Path:
    """PPTX com slide contendo tabela 2×3."""
    prs = Presentation()
    layout = prs.slide_layouts[5]  # em branco
    slide = prs.slides.add_slide(layout)

    rows, cols = 3, 2
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(3)).table
    dados = [["Nome", "Valor"], ["Alpha", "100"], ["Beta", "200"]]
    for r, row_data in enumerate(dados):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    path = tmp_path / "tabela.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def pptx_vazio(tmp_path) -> Path:
    """PPTX com slide sem conteúdo de texto."""
    prs = Presentation()
    layout = prs.slide_layouts[5]  # em branco
    prs.slides.add_slide(layout)
    path = tmp_path / "vazio.pptx"
    prs.save(str(path))
    return path


# ── Testes básicos ───────────────────────────────────────────────────────────

def test_pptx_dois_slides(pptx_basico):
    """Dois slides → dois headings ## e dois ### (títulos)."""
    md = pptx_to_md(pptx_basico)
    assert "## Slide 1" in md
    assert "## Slide 2" in md
    assert "### Slide Um" in md
    assert "### Slide Dois" in md


def test_pptx_conteudo_corpo(pptx_basico):
    """Conteúdo do corpo aparece no output."""
    md = pptx_to_md(pptx_basico)
    assert "Conteúdo do slide um." in md
    assert "Conteúdo do slide dois." in md


def test_pptx_tabela(pptx_com_tabela):
    """Tabela → formato Markdown com header + separador + linhas."""
    md = pptx_to_md(pptx_com_tabela)
    assert "| Nome | Valor |" in md
    assert "| --- | --- |" in md
    assert "| Alpha | 100 |" in md
    assert "| Beta | 200 |" in md


def test_pptx_vazio_nao_falha(pptx_vazio):
    """Slide sem texto não levanta exceção."""
    md = pptx_to_md(pptx_vazio)
    assert isinstance(md, str)


def test_pptx_arquivo_inexistente(tmp_path):
    """FileNotFoundError para arquivo inexistente."""
    with pytest.raises(FileNotFoundError):
        pptx_to_md(tmp_path / "nao_existe.pptx")


def test_pptx_extensao_errada(tmp_path):
    """ValueError para extensão incorreta."""
    f = tmp_path / "arquivo.pdf"
    f.write_bytes(b"%PDF-1.4")
    with pytest.raises(ValueError):
        pptx_to_md(f)


def test_pptx_corrompido(tmp_path):
    """RuntimeError para PPTX inválido."""
    f = tmp_path / "corrompido.pptx"
    f.write_bytes(b"nao_e_pptx")
    with pytest.raises(RuntimeError):
        pptx_to_md(f)


# ── Testes do helper _tabela_para_md ─────────────────────────────────────────

def test_tabela_para_md_pipe_em_celula(pptx_com_tabela):
    """Pipes dentro de células são escapados."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "A|B"
    table.cell(0, 1).text = "C"
    table.cell(1, 0).text = "D"
    table.cell(1, 1).text = "E"

    md = _tabela_para_md(table)
    assert "A\\|B" in md
