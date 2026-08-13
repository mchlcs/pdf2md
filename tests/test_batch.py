"""Testes para core/batch.py."""

from pathlib import Path

import fitz
import pytest
from PIL import Image, ImageDraw

from core.batch import StatusArquivo, batch_convert
from core.utils import ModoImagem


def test_batch_arquivo_unico(fixture_texto_simples, tmp_path):
    """Arquivo único → 1 MD."""
    resultados = batch_convert(
        origem=fixture_texto_simples,
        destino=tmp_path / "saida",
        workers=1,
    )
    assert len(resultados) == 1
    assert resultados[0].status == StatusArquivo.CONCLUIDO
    assert resultados[0].destino is not None
    assert resultados[0].destino.exists()


def test_batch_diretorio(tmp_path):
    """Dir com 3 PDFs → 3 MDs."""
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    for i in range(3):
        doc = fitz.open()
        p = doc.new_page(width=595, height=842)
        p.insert_text((50, 50), f"PDF {i}")
        doc.save(str(entrada / f"doc{i}.pdf"))
        doc.close()

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida",
        workers=2,
    )
    assert len(resultados) == 3
    assert all(r.status == StatusArquivo.CONCLUIDO for r in resultados)


def test_batch_mix_pdf_imagem(tmp_path, fixture_img_texto, mock_tesseract):
    """2 PDFs + 2 PNGs → 4 MDs."""
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    for i in range(2):
        doc = fitz.open()
        p = doc.new_page(width=595, height=842)
        p.insert_text((50, 50), f"PDF {i}")
        doc.save(str(entrada / f"doc{i}.pdf"))
        doc.close()

    # Copia imagem fixture
    import shutil
    shutil.copy(fixture_img_texto, entrada / "img.png")
    # Cria segunda imagem
    img = Image.new("RGB", (400, 200), color="white")
    draw = ImageDraw.Draw(img)
    draw.text((10, 80), "Segunda imagem", fill="black")
    img.save(entrada / "img2.jpg")

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida",
        workers=2,
    )
    assert len(resultados) == 4
    assert all(r.status == StatusArquivo.CONCLUIDO for r in resultados)


def test_batch_nao_sobrescreve(tmp_path):
    """sobrescrever=False → pula existentes."""
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    doc = fitz.open()
    p = doc.new_page(width=595, height=842)
    p.insert_text((50, 50), "Texto")
    doc.save(str(entrada / "doc.pdf"))
    doc.close()

    saida = tmp_path / "saida"
    saida.mkdir()
    (saida / "doc.md").write_text("já existe", encoding="utf-8")

    resultados = batch_convert(
        origem=entrada / "doc.pdf",
        destino=saida,
        workers=1,
        sobrescrever=False,
    )
    assert resultados[0].status == StatusArquivo.IGNORADO
    # IGNORADO distingue "pulado por já existir" de "convertido com sucesso"
    # Verificamos que o conteúdo não foi alterado
    assert (saida / "doc.md").read_text(encoding="utf-8") == "já existe"


def test_batch_vault_salva_direto(tmp_path):
    """vault definido → salva diretamente em vault/ (sem subpasta _inbox)."""
    vault = tmp_path / "vault"
    vault.mkdir()
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    doc = fitz.open()
    p = doc.new_page(width=595, height=842)
    p.insert_text((50, 50), "Vault test")
    doc.save(str(entrada / "doc.pdf"))
    doc.close()

    resultados = batch_convert(
        origem=entrada / "doc.pdf",
        destino=tmp_path / "saida",
        vault=vault,
        workers=1,
    )
    assert resultados[0].status == StatusArquivo.CONCLUIDO
    assert (vault / "doc.md").exists()


def test_batch_obsidian_frontmatter(tmp_path):
    """obsidian=True → MD tem frontmatter."""
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    doc = fitz.open()
    p = doc.new_page(width=595, height=842)
    p.insert_text((50, 50), "Frontmatter test")
    doc.save(str(entrada / "doc.pdf"))
    doc.close()

    resultados = batch_convert(
        origem=entrada / "doc.pdf",
        destino=tmp_path / "saida",
        obsidian=True,
        workers=1,
    )
    assert resultados[0].status == StatusArquivo.CONCLUIDO
    md = (tmp_path / "saida" / "doc.md").read_text(encoding="utf-8")
    assert "---" in md
    assert "title:" in md
    assert "pdf2md" in md


def test_batch_extensao_ignorada(tmp_path):
    """.txt → StatusArquivo.IGNORADO."""
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    (entrada / "leitura.txt").write_text("não suportado", encoding="utf-8")

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida",
        workers=1,
    )
    assert len(resultados) == 1
    assert resultados[0].status == StatusArquivo.IGNORADO


def test_batch_vault_invalido(tmp_path):
    """NotADirectoryError se vault não é diretório."""
    with pytest.raises(NotADirectoryError):
        batch_convert(
            origem=tmp_path,
            destino=tmp_path / "saida",
            vault=tmp_path / "nao_e_dir.txt",
            workers=1,
        )


def test_batch_path_traversal(tmp_path):
    """ValueError para ../../etc/passwd.pdf."""
    with pytest.raises(ValueError, match="traversal"):
        batch_convert(
            origem=tmp_path / ".." / ".." / "etc" / "passwd.pdf",
            destino=tmp_path / "saida",
            workers=1,
        )


def test_batch_sem_duracao_por_arquivo(tmp_path):
    """ResultadoArquivo não tem mais campo duracao — tempo é só total (CLI)."""
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    doc = fitz.open()
    p = doc.new_page(width=595, height=842)
    p.insert_text((50, 50), "Texto suficiente para validar a conversao do PDF.")
    doc.save(str(entrada / "doc.pdf"))
    doc.close()

    resultados = batch_convert(origem=entrada, destino=tmp_path / "saida", workers=1)
    assert len(resultados) == 1
    assert resultados[0].status == StatusArquivo.CONCLUIDO
    assert not hasattr(resultados[0], "duracao")  # campo removido


def test_batch_colisao_stem_nao_perde_dados(tmp_path):
    """report.pdf + report.docx → 2 MDs distintos, sem sobrescrita (data race).

    Regressão: antes ambos geravam report.md → write_text concorrente sob
    ThreadPoolExecutor causava perda silenciosa de um dos arquivos.
    """
    from docx import Document

    entrada = tmp_path / "entrada"
    entrada.mkdir()

    # PDF com texto nativo (>50 chars ASCII → caminho pymupdf4llm, sem OCR)
    doc = fitz.open()
    p = doc.new_page(width=595, height=842)
    p.insert_text((50, 50), "Conteudo do PDF de relatorio para validacao do batch converter.")
    doc.save(str(entrada / "report.pdf"))
    doc.close()

    # DOCX com o MESMO stem
    docx = Document()
    docx.add_paragraph("Conteudo do DOCX de relatorio para validacao.")
    docx.save(str(entrada / "report.docx"))

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida",
        workers=2,
    )

    concluidos = [r for r in resultados if r.status == StatusArquivo.CONCLUIDO]
    assert len(concluidos) == 2
    # Destinos precisam ser distintos — sem colisão de nome
    destinos = {r.destino for r in concluidos}
    assert len(destinos) == 2
    for d in destinos:
        assert d is not None and d.exists()
        assert d.read_text(encoding="utf-8").strip() != ""


# ── Feature --imagens (T4): propagação e assets por arquivo ──────────────────

def _pdf_com_imagem_e_texto(tmp_path, nome: str) -> fitz.Document:
    """Cria PDF com texto nativo + 1 imagem embutida (cor distinta por nome)."""
    path = tmp_path / f"{nome}.pdf"
    doc = fitz.open()
    p = doc.new_page(width=595, height=842)
    p.insert_text(
        (50, 150),
        "Documento com texto nativo e imagem embutida para testes de assets do batch.",
    )
    cor = ((hash(nome) % 200) + 20, 40, 90)
    img = Image.new("RGB", (40, 40), color=cor)
    img_path = tmp_path / f"{nome}_img.png"
    img.save(img_path)
    p.insert_image(fitz.Rect(50, 50, 90, 90), filename=str(img_path))
    doc.save(str(path))
    doc.close()
    return path


def test_batch_imagens_extrair_assets_por_arquivo(tmp_path):
    """--imagens extrair: assets por arquivo em <stem>_assets/, sem colisão."""
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    _pdf_com_imagem_e_texto(entrada, "alfa")
    _pdf_com_imagem_e_texto(entrada, "beta")

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida",
        workers=2,
        modo_imagem=ModoImagem.extrair,
    )

    assert all(r.status == StatusArquivo.CONCLUIDO for r in resultados)
    for stem in ("alfa", "beta"):
        assets = tmp_path / "saida" / f"{stem}_assets"
        assert (assets / "img_p001_0.png").exists()
        md = (tmp_path / "saida" / f"{stem}.md").read_text(encoding="utf-8")
        assert f"{stem}_assets/img_p001_0.png" in md


def test_batch_imagens_assets_compartilhado_com_prefixo(tmp_path):
    """--assets-dir compartilhado: prefixo por stem evita colisão (T4)."""
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    _pdf_com_imagem_e_texto(entrada, "alfa")
    _pdf_com_imagem_e_texto(entrada, "beta")
    assets_globais = tmp_path / "assets_globais"

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida",
        workers=2,
        modo_imagem=ModoImagem.extrair,
        assets_dir=assets_globais,
    )

    assert all(r.status == StatusArquivo.CONCLUIDO for r in resultados)
    assert (assets_globais / "alfa__img_p001_0.png").exists()
    assert (assets_globais / "beta__img_p001_0.png").exists()
    md_alfa = (tmp_path / "saida" / "alfa.md").read_text(encoding="utf-8")
    assert "alfa__img_p001_0.png" in md_alfa


def test_batch_imagens_fora_do_pdf_docx_emite_aviso(tmp_path):
    """--imagens em PPTX → aviso, não erro (semântica PDF/DOCX-only)."""
    from pptx import Presentation

    entrada = tmp_path / "entrada"
    entrada.mkdir()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide de teste"
    prs.save(str(entrada / "apresentacao.pptx"))

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida",
        workers=1,
        modo_imagem=ModoImagem.extrair,
    )

    assert len(resultados) == 1
    assert resultados[0].status == StatusArquivo.CONCLUIDO
    assert any("só se aplica a PDF e DOCX" in aviso for aviso in resultados[0].avisos)


def test_batch_imagens_obsidian_attachments(tmp_path):
    """--obsidian + --imagens: assets em vault/attachments com wikilink."""
    entrada = tmp_path / "entrada"
    entrada.mkdir()
    _pdf_com_imagem_e_texto(entrada, "nota")
    vault = tmp_path / "vault"
    vault.mkdir()

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida_ignorada",
        vault=vault,
        workers=1,
        modo_imagem=ModoImagem.extrair,
    )

    assert resultados[0].status == StatusArquivo.CONCLUIDO
    attachments = vault / "attachments"
    assert (attachments / "nota__img_p001_0.png").exists()
    md = (vault / "nota.md").read_text(encoding="utf-8")
    assert "![[nota__img_p001_0.png]]" in md


def test_batch_docx_imagens_extrair(tmp_path):
    """DOCX + --imagens extrair: assets por arquivo, sem aviso de formato."""
    import tempfile

    from docx import Document
    from PIL import Image

    entrada = tmp_path / "entrada"
    entrada.mkdir()
    docx = Document()
    docx.add_paragraph("Relatorio com figura.")
    p = docx.add_paragraph()
    run = p.add_run()
    img = Image.new("RGB", (30, 30), color=(30, 120, 200))
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        img.save(tmp.name)
        run.add_picture(tmp.name)
    docx.save(str(entrada / "relatorio.docx"))

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida",
        workers=1,
        modo_imagem=ModoImagem.extrair,
    )

    assert resultados[0].status == StatusArquivo.CONCLUIDO
    assert not any("só se aplica a PDF e DOCX" in aviso for aviso in resultados[0].avisos)
    assets = tmp_path / "saida" / "relatorio_assets"
    assert (assets / "img_0001.png").exists()
    md = (tmp_path / "saida" / "relatorio.md").read_text(encoding="utf-8")
    assert "relatorio_assets/img_0001.png" in md
    assert "base64" not in md


def test_batch_docx_obsidian_wikilink(tmp_path):
    """DOCX + --obsidian: wikilink e assets em vault/attachments."""
    import tempfile

    from docx import Document
    from PIL import Image

    entrada = tmp_path / "entrada"
    entrada.mkdir()
    docx = Document()
    docx.add_paragraph("Nota com figura.")
    p = docx.add_paragraph()
    run = p.add_run()
    img = Image.new("RGB", (30, 30), color=(90, 200, 40))
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        img.save(tmp.name)
        run.add_picture(tmp.name)
    docx.save(str(entrada / "nota.docx"))
    vault = tmp_path / "vault"
    vault.mkdir()

    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "ignorado",
        vault=vault,
        workers=1,
        modo_imagem=ModoImagem.extrair,
    )

    assert resultados[0].status == StatusArquivo.CONCLUIDO
    assert (vault / "attachments" / "nota__img_0001.png").exists()
    md = (vault / "nota.md").read_text(encoding="utf-8")
    assert "![[nota__img_0001.png]]" in md


def test_batch_imagens_stem_colidido_nao_sobrescreve(tmp_path):
    """Mesmo stem com extensões diferentes + --assets-dir: prefixos distintos.

    Review Spec-3: o prefixo agora deriva do stem do .md destino (já único
    via _nome_destino_unico), não do stem da origem — antes, `rel.pdf` +
    `rel.docx` geravam prefixos idênticos e sobrescreviam assets.
    """
    import tempfile

    from docx import Document
    from PIL import Image

    entrada = tmp_path / "entrada"
    entrada.mkdir()

    # PDF "rel" com imagem
    _pdf_com_imagem_e_texto(entrada, "rel")
    # DOCX "rel" com imagem
    docx = Document()
    docx.add_paragraph("Relatorio docx.")
    p = docx.add_paragraph()
    run = p.add_run()
    img = Image.new("RGB", (30, 30), color=(150, 40, 40))
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        img.save(tmp.name)
        run.add_picture(tmp.name)
    docx.save(str(entrada / "rel.docx"))

    assets_globais = tmp_path / "assets_globais"
    resultados = batch_convert(
        origem=entrada,
        destino=tmp_path / "saida",
        workers=2,
        modo_imagem=ModoImagem.extrair,
        assets_dir=assets_globais,
    )

    assert all(r.status == StatusArquivo.CONCLUIDO for r in resultados)
    arquivos = sorted(f.name for f in assets_globais.iterdir())
    assert len(arquivos) == 2  # nenhum asset foi sobrescrito
    # Ordenação: rel.docx → rel.md (prefixo rel__), rel.pdf → rel-pdf.md
    assert arquivos == ["rel-pdf__img_p001_0.png", "rel__img_0001.png"]


# ── TAREFA 5: exceção real no worker, colisão 3+, subdiretórios ─────────────

def test_batch_pdf_corrompido_erro_sanitizado_sem_interromper(tmp_path):
    """PDF corrompido no diretório → ERRO sanitizado; o batch continua.

    Regressão de robustez: a exceção real do worker não pode derrubar os
    demais arquivos nem vazar paths absolutos na mensagem de erro.
    """
    entrada = tmp_path / "entrada"
    entrada.mkdir()

    # PDF válido
    doc = fitz.open()
    p = doc.new_page(width=595, height=842)
    p.insert_text((50, 50), "PDF valido do lote.")
    doc.save(str(entrada / "bom.pdf"))
    doc.close()

    # PDF corrompido (não abre no PyMuPDF)
    (entrada / "quebrado.pdf").write_bytes(b"%PDF-1.4\n% lixo nao estrutura valida")

    resultados = batch_convert(origem=entrada, destino=tmp_path / "saida", workers=2)

    assert len(resultados) == 2
    por_status = {r.status: r for r in resultados}
    assert por_status[StatusArquivo.CONCLUIDO].destino.exists()
    erro = por_status[StatusArquivo.ERRO]
    assert erro.erro is not None and erro.erro != ""
    assert "Traceback" not in erro.erro
    assert str(tmp_path) not in erro.erro  # path sanitizado (CWE-209)


def test_processar_arquivo_excecao_vira_erro_sanitizado(tmp_path):
    """Worker direto: exceção real → dict de ERRO, sem propagar."""
    from core.batch import _processar_arquivo

    origem = tmp_path / "quebrado.pdf"
    origem.write_bytes(b"%PDF-1.4\n% lixo")

    resultado = _processar_arquivo(str(origem), str(tmp_path / "saida" / "quebrado.md"), False, False)

    assert resultado["status"] == "erro"
    assert resultado["destino"] is None
    assert resultado["erro"] is not None
    assert str(tmp_path) not in resultado["erro"]  # sanitizado


def test_batch_colisao_stem_tres_arquivos(tmp_path):
    """3 arquivos com o mesmo stem (pdf/docx/pptx) → 3 MDs distintos.

    Estende a regressão de data race para 3+ formatos: cada arquivo ganha
    nome único via sufixo de extensão, sem sobrescrita silenciosa.
    """
    from docx import Document
    from pptx import Presentation

    entrada = tmp_path / "entrada"
    entrada.mkdir()

    doc = fitz.open()
    p = doc.new_page(width=595, height=842)
    p.insert_text((50, 50), "Conteudo do PDF de relatorio para validacao do batch converter.")
    doc.save(str(entrada / "rel.pdf"))
    doc.close()

    docx = Document()
    docx.add_paragraph("Conteudo do DOCX de relatorio para validacao.")
    docx.save(str(entrada / "rel.docx"))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Relatorio PPTX"
    prs.save(str(entrada / "rel.pptx"))

    resultados = batch_convert(origem=entrada, destino=tmp_path / "saida", workers=3)

    concluidos = [r for r in resultados if r.status == StatusArquivo.CONCLUIDO]
    assert len(concluidos) == 3
    destinos = {r.destino for r in concluidos}
    assert len(destinos) == 3  # nenhuma colisão de nome
    for d in destinos:
        assert d is not None and d.exists()
        assert d.read_text(encoding="utf-8").strip() != ""


def test_nome_destino_unico_contador_2():
    """Colisão no candidato com sufixo → contador -2, -3, ..."""
    from core.batch import _nome_destino_unico

    usados = {"rel.md", "rel-docx.md"}
    assert _nome_destino_unico(Path("rel.docx"), usados) == "rel-docx-2.md"
    usados.add("rel-docx-2.md")
    assert _nome_destino_unico(Path("rel.docx"), usados) == "rel-docx-3.md"


def test_batch_subdiretorio_ignorado(tmp_path):
    """Subdiretório dentro do lote é ignorado (varredura não recursiva)."""
    entrada = tmp_path / "entrada"
    sub = entrada / "subpasta"
    sub.mkdir(parents=True)

    doc = fitz.open()
    p = doc.new_page(width=595, height=842)
    p.insert_text((50, 50), "PDF da raiz.")
    doc.save(str(entrada / "raiz.pdf"))
    doc.close()

    doc2 = fitz.open()
    p2 = doc2.new_page(width=595, height=842)
    p2.insert_text((50, 50), "PDF dentro da subpasta.")
    doc2.save(str(sub / "aninhado.pdf"))
    doc2.close()

    resultados = batch_convert(origem=entrada, destino=tmp_path / "saida", workers=1)

    assert len(resultados) == 1
    assert resultados[0].status == StatusArquivo.CONCLUIDO
    assert resultados[0].origem.name == "raiz.pdf"
    assert not (tmp_path / "saida" / "aninhado.md").exists()


def test_batch_diretorio_apenas_subpasta_sem_resultados(tmp_path):
    """Diretório só com subpasta → nenhum resultado, sem crash."""
    entrada = tmp_path / "entrada"
    (entrada / "subpasta").mkdir(parents=True)
    (entrada / "subpasta" / "x.pdf").write_bytes(b"%PDF-1.4 dummy")

    resultados = batch_convert(origem=entrada, destino=tmp_path / "saida", workers=1)

    assert resultados == []
