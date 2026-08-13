"""
Converte arquivos .pptx em Markdown.
Estratégia: python-pptx → extrai título + corpo + tabelas por slide.
"""
from pathlib import Path
from typing import Any

from core.utils import _sanitizar_celula_md, _validar_existencia


def pptx_to_md(path: Path) -> str:
    """
    Converte arquivo .pptx em Markdown.

    Estratégia:
    - Cada slide → seção (## Slide N)
    - Placeholder título (idx=0) → heading (### texto)
    - Corpo → parágrafos não-vazios preservados
    - Tabelas → tabela Markdown
    - Imagens → omitidas (sem OCR; usar image_converter separadamente se necessário)

    Args:
        path: Caminho absoluto para o arquivo .pptx. Deve existir.

    Returns:
        String Markdown com conteúdo extraído de todos os slides.

    Raises:
        FileNotFoundError: Se path não existe.
        ValueError: Se path não é arquivo .pptx válido.
        RuntimeError: Se python-pptx falha ao abrir o arquivo.
    """
    from pptx import Presentation  # lazy import — evita custo de import no startup

    _validar_existencia(path)
    if path.suffix.lower() != ".pptx":
        raise ValueError("Arquivo não é PPTX válido")

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise RuntimeError(
            "Falha ao abrir PPTX — arquivo corrompido ou formato inválido"
        ) from exc

    partes: list[str] = []

    for num_slide, slide in enumerate(prs.slides, 1):
        partes.append(f"## Slide {num_slide}")
        for shape in slide.shapes:
            _processar_shape(shape, partes)

    return "\n\n".join(partes)


def _processar_shape(shape: Any, partes: list[str]) -> None:
    """Extrai conteúdo de um shape (tabela, título ou corpo) para a lista de partes."""
    if shape.has_table:
        md_tabela = _tabela_para_md(shape.table)
        if md_tabela:
            partes.append(md_tabela)
        return

    if not shape.has_text_frame:
        return

    texto_shape = shape.text_frame.text.strip()
    if not texto_shape:
        return

    eh_titulo = (
        hasattr(shape, "placeholder_format")
        and shape.placeholder_format is not None
        and shape.placeholder_format.idx == 0
    )

    if eh_titulo:
        partes.append(f"### {texto_shape}")
    else:
        for para in shape.text_frame.paragraphs:
            linha = para.text.strip()
            if linha:
                partes.append(linha)


def _tabela_para_md(table: Any) -> str:
    """Converte python-pptx Table em tabela Markdown."""
    if not table.rows:
        return ""

    linhas: list[str] = []
    for i, row in enumerate(table.rows):
        celulas = [_sanitizar_celula_md(cell.text) for cell in row.cells]
        linhas.append("| " + " | ".join(celulas) + " |")
        if i == 0:
            linhas.append("| " + " | ".join(["---"] * len(celulas)) + " |")

    return "\n".join(linhas)
